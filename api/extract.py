
import io
import re
from pypdf import PdfReader
from pptx import Presentation
from docx import Document
import openpyxl
from flask import Flask, request, jsonify

app = Flask(__name__)

# Mengatur batas maksimal konten agar Flask menerima file hingga 25MB
app.config['MAX_CONTENT_LENGTH'] = 25 * 1024 * 1024

def clean_text(text):
    # Membersihkan karakter aneh dan spasi berlebih
    text = re.sub(r'([A-Z])\s(?=[a-z])', r'\1', text) 
    return " ".join(text.split())

def extract_metadata_heuristics(full_text, filename):
    metadata = {
        "title": filename.rsplit('.', 1)[0].replace("_", " "),
        "authors": [],
        "year": "",
        "publisher": "",
        "keywords": [],
        "category": "Original Research",
        "type": "Literature"
    }

    # Ekstraksi Tahun
    year_match = re.search(r'\b(19|20)\d{2}\b', full_text[:5000])
    if year_match:
        metadata["year"] = year_match.group(0)

    # Heuristik Penerbit Sederhana
    publishers = ["Elsevier", "Springer", "IEEE", "MDPI", "Nature", "Science", "Wiley", "Taylor & Francis", "ACM", "Frontiers", "Sage"]
    for pub in publishers:
        if pub.lower() in full_text[:10000].lower():
            metadata["publisher"] = pub
            break

    return metadata

@app.route('/api/extract', methods=['POST'])
def extract():
    try:
        if 'file' not in request.files:
            return jsonify({"status": "error", "message": "No file part"}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({"status": "error", "message": "No file selected"}), 400

        filename_lower = file.filename.lower()
        
        # Blokir file Audio dan Video secara eksplisit
        audio_video_ext = ('.mp3', '.wav', '.ogg', '.m4a', '.flac', '.aac', '.mp4', '.avi', '.mov', '.mkv', '.wmv', '.flv', '.webm')
        if filename_lower.endswith(audio_video_ext):
            return jsonify({"status": "error", "message": "Audio and video files are not supported. Please upload a document."}), 400

        # Handling file legacy (.doc, .xls, .ppt)
        legacy_ext = ('.doc', '.xls', '.ppt')
        if filename_lower.endswith(legacy_ext):
             return jsonify({
                "status": "error", 
                "message": f"Legacy format {filename_lower.split('.')[-1]} detected. Please save/convert as modern format (e.g. .docx, .xlsx, .pptx) before uploading."
            }), 422

        file_bytes = file.read()
        f = io.BytesIO(file_bytes)
        full_text = ""

        # Ekstraksi berdasarkan format file
        if filename_lower.endswith('.pdf'):
            try:
                reader = PdfReader(f)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        full_text += page_text + "\n"
            except Exception as e:
                return jsonify({"status": "error", "message": f"PDF Extraction Error: {str(e)}"}), 422
                
        elif filename_lower.endswith('.docx'):
            try:
                doc = Document(f)
                full_text = "\n".join([para.text for para in doc.paragraphs])
            except Exception as e:
                return jsonify({"status": "error", "message": f"Word (.docx) Error: {str(e)}"}), 422

        elif filename_lower.endswith('.pptx'):
            try:
                prs = Presentation(f)
                for slide in prs.slides:
                    # Ambil teks dari shape/box di slide
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            full_text += shape.text + "\n"
                    # Ambil teks dari catatan pembicara (notes)
                    if slide.has_notes_slide:
                        notes_text = slide.notes_slide.notes_text_frame.text
                        if notes_text:
                            full_text += notes_text + "\n"
            except Exception as e:
                return jsonify({"status": "error", "message": f"PowerPoint (.pptx) Error: {str(e)}"}), 422

        elif filename_lower.endswith('.xlsx'):
            try:
                wb = openpyxl.load_workbook(f, data_only=True)
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        full_text += " ".join([str(cell) for cell in row if cell is not None]) + "\n"
            except Exception as e:
                return jsonify({"status": "error", "message": f"Excel (.xlsx) Error: {str(e)}"}), 422

        elif filename_lower.endswith(('.txt', '.md', '.csv')):
            try:
                full_text = file_bytes.decode('utf-8', errors='ignore')
            except Exception as e:
                return jsonify({"status": "error", "message": f"Text File Error: {str(e)}"}), 422
        else:
            return jsonify({"status": "error", "message": "Unsupported file format. Please upload a common document format."}), 400

        if not full_text.strip():
            return jsonify({"status": "error", "message": "Document is empty or text could not be extracted."}), 422

        # 1. Batasi total teks (200.000 karakter)
        limit_total = 200000
        limited_text = full_text[:limit_total]

        # 2. Heuristik metadata
        metadata = extract_metadata_heuristics(limited_text, file.filename)
        
        # 3. Snippet untuk AI (Groq/Gemini) - 7500 karakter
        ai_snippet = limited_text[:7500]
        
        # 4. Split teks ke dalam 10 chunks (masing-masing 20.000 karakter)
        chunk_size = 20000
        chunks = [limited_text[i:i+chunk_size] for i in range(0, len(limited_text), chunk_size)][:10]

        return jsonify({
            "status": "success",
            "data": {
                **metadata,
                "aiSnippet": ai_snippet,
                "chunks": chunks
            }
        })
    except Exception as e:
        return jsonify({"status": "error", "message": f"Internal Server Error: {str(e)}"}), 500

if __name__ == '__main__':
    app.run(port=5000)
